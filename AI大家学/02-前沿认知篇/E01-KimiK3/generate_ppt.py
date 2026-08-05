"""
Kimi K3 视频演示文稿 v5 — 胡桃色温馨风格 + 动画效果
python-pptx 生成，直接注入动画 XML
v5 更新：新增开源三件套、商业数据、本地部署番外（共25页）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ========== 配色 ==========
C = {
    'bg':      RGBColor(0xF5,0xF0,0xE8),
    'card':    RGBColor(0xFA,0xF7,0xF2),
    'primary': RGBColor(0x8B,0x6F,0x47),
    'accent':  RGBColor(0xE8,0xA8,0x7C),
    'teal':    RGBColor(0x7B,0xA7,0xA7),
    'text':    RGBColor(0x5D,0x4E,0x37),
    'subtext': RGBColor(0x8B,0x7D,0x6B),
    'border':  RGBColor(0xED,0xE5,0xD5),
    'white':   RGBColor(0xFF,0xFF,0xFF),
    'red':     RGBColor(0xC2,0x5B,0x56),
    'green':   RGBColor(0x7B,0xA7,0x7B),
    'purple':  RGBColor(0x8B,0x7B,0xA7),
    'blue':    RGBColor(0x5B,0x7B,0xA7),
    'dark_card': RGBColor(0xF0,0xEB,0xE0),
}
FONT_CN = '微软雅黑'
FONT_EN = 'Georgia'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # blank layout

# ========== 辅助函数 ==========
def set_bg(slide, color=C['bg']):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=14, bold=False,
             color=C['text'], font=FONT_CN, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox

def add_text_multi(slide, left, top, width, height, segments, align=PP_ALIGN.LEFT):
    """segments: list of (text, size, bold, color, font)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (text, size, bold, color, font) in enumerate(segments):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return txBox

def add_text_multiline(slide, left, top, width, height, lines, size=13,
                        color=C['text'], font=FONT_CN, align=PP_ALIGN.LEFT, line_spacing=None):
    """lines: list of str, each is a paragraph"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = font
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_ellipse(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_footer(slide, text):
    add_text(slide, 0, 6.8, 13.33, 0.3, text, size=8, color=C['subtext'], align=PP_ALIGN.CENTER)

def add_section_page(slide, num, title, subtitle):
    set_bg(slide)
    add_text(slide, 0.6, 0.8, 1.5, 1.2, num, size=60, bold=True, color=C['border'], font=FONT_EN)
    add_text(slide, 2.2, 1.0, 8, 1.0, title, size=36, bold=True, color=C['primary'])
    add_text(slide, 2.2, 2.0, 8, 0.5, subtitle, size=16, color=C['subtext'])
    add_rect(slide, 2.2, 2.6, 3, 0.03, fill_color=C['accent'])

# ========== 动画系统 ==========
_anim_id_counter = [0]

def _next_id():
    _anim_id_counter[0] += 1
    return _anim_id_counter[0]

def _build_entrance_anim(shape, anim_type='fade', delay_ms=0, dur_ms=500, node_type='afterEffect'):
    sp_id = shape.shape_id
    preset_ids = {
        'fade': 10, 'fly_left': 2, 'fly_right': 2, 'fly_top': 2, 'fly_bottom': 2,
        'wipe_right': 22, 'wipe_bottom': 22, 'zoom': 53,
    }
    preset_subtypes = {
        'fade': '0', 'fly_left': '4', 'fly_right': '8', 'fly_top': '2', 'fly_bottom': '1',
        'wipe_right': '4', 'wipe_bottom': '2', 'zoom': '0',
    }
    preset_id = preset_ids.get(anim_type, 10)
    sub_type = preset_subtypes.get(anim_type, '0')

    par = etree.SubElement(etree.Element('dummy'), qn('p:par'))
    ctn_par = etree.SubElement(par, qn('p:cTn'))
    ctn_par.set('id', str(_next_id()))
    ctn_par.set('presetID', str(preset_id))
    ctn_par.set('presetClass', 'entr')
    ctn_par.set('presetSubtype', sub_type)
    ctn_par.set('fill', 'hold')
    ctn_par.set('nodeType', node_type)

    stCond = etree.SubElement(ctn_par, qn('p:stCondLst'))
    cond = etree.SubElement(stCond, qn('p:cond'))
    cond.set('delay', str(delay_ms))

    childTn = etree.SubElement(ctn_par, qn('p:childTnLst'))
    set_el = etree.SubElement(childTn, qn('p:set'))
    set_bhvr = etree.SubElement(set_el, qn('p:cBhvr'))
    set_ctn = etree.SubElement(set_bhvr, qn('p:cTn'))
    set_ctn.set('id', str(_next_id()))
    set_ctn.set('dur', '1')
    set_ctn.set('fill', 'hold')
    stCond2 = etree.SubElement(set_ctn, qn('p:stCondLst'))
    cond2 = etree.SubElement(stCond2, qn('p:cond'))
    cond2.set('delay', '0')
    tgtEl = etree.SubElement(set_bhvr, qn('p:tgtEl'))
    spTgt = etree.SubElement(tgtEl, qn('p:spTgt'))
    spTgt.set('spid', str(sp_id))
    attrNameLst = etree.SubElement(set_bhvr, qn('p:attrNameLst'))
    attrName = etree.SubElement(attrNameLst, qn('p:attrName'))
    attrName.text = 'style.visibility'
    to_el = etree.SubElement(set_el, qn('p:to'))
    val = etree.SubElement(to_el, qn('p:strVal'))
    val.set('val', 'visible')

    anim_effect = etree.SubElement(childTn, qn('p:animEffect'))
    anim_effect.set('transition', 'in')
    anim_effect.set('filter', f'preset={preset_id}')
    ae_bhvr = etree.SubElement(anim_effect, qn('p:cBhvr'))
    ae_ctn = etree.SubElement(ae_bhvr, qn('p:cTn'))
    ae_ctn.set('id', str(_next_id()))
    ae_ctn.set('dur', str(dur_ms))
    ae_tgt = etree.SubElement(ae_bhvr, qn('p:tgtEl'))
    ae_sp = etree.SubElement(ae_tgt, qn('p:spTgt'))
    ae_sp.set('spid', str(sp_id))

    return par


def add_slide_transition(slide, trans_type='fade', speed_ms=700):
    timing = slide._element.find(qn('p:transition'))
    if timing is not None:
        slide._element.remove(timing)
    trans = etree.SubElement(slide._element, qn('p:transition'))
    trans.set('spd', 'med')
    trans.set('advClick', '1')
    trans_types = {'fade': 'fade', 'push': 'push', 'wipe': 'wipe', 'cover': 'cover', 'dissolve': 'dissolve'}
    tag = trans_types.get(trans_type, 'fade')
    el = etree.SubElement(trans, qn(f'p:{tag}'))
    if tag in ('push', 'wipe', 'cover'):
        el.set('dir', 'l' if trans_type == 'push' else 'd')


def animate(slide, shape, anim_type='fade', delay_s=0, dur_s=0.5):
    if not hasattr(slide, '_anim_queue'):
        slide._anim_queue = []
    delay_ms = int(delay_s * 1000)
    dur_ms = int(dur_s * 1000)
    slide._anim_queue.append((shape, anim_type, delay_ms, dur_ms))


def flush_animations(slide):
    if not hasattr(slide, '_anim_queue') or not slide._anim_queue:
        return
    queue = slide._anim_queue
    slide._anim_queue = []

    existing = slide._element.find(qn('p:timing'))
    if existing is not None:
        slide._element.remove(existing)

    timing = etree.SubElement(slide._element, qn('p:timing'))
    tnLst = etree.SubElement(timing, qn('p:tnLst'))
    root_par = etree.SubElement(tnLst, qn('p:par'))
    root_ctn = etree.SubElement(root_par, qn('p:cTn'))
    root_ctn.set('id', str(_next_id()))
    root_ctn.set('dur', 'indefinite')
    root_ctn.set('restart', 'never')
    root_ctn.set('nodeType', 'tmRoot')

    root_child = etree.SubElement(root_ctn, qn('p:childTnLst'))
    seq = etree.SubElement(root_child, qn('p:seq'))
    seq.set('concurrent', '1')
    seq.set('nextAc', 'seek')
    seq_ctn = etree.SubElement(seq, qn('p:cTn'))
    seq_ctn.set('id', str(_next_id()))
    seq_ctn.set('dur', 'indefinite')
    seq_ctn.set('nodeType', 'mainSeq')

    seq_child = etree.SubElement(seq_ctn, qn('p:childTnLst'))
    group = etree.SubElement(seq_child, qn('p:par'))
    group_ctn = etree.SubElement(group, qn('p:cTn'))
    group_ctn.set('id', str(_next_id()))
    group_ctn.set('fill', 'hold')
    group_stCond = etree.SubElement(group_ctn, qn('p:stCondLst'))
    group_cond = etree.SubElement(group_stCond, qn('p:cond'))
    group_cond.set('delay', '0')
    group_child = etree.SubElement(group_ctn, qn('p:childTnLst'))

    prev_end_ms = 0
    for i, (shape, anim_type, delay_ms, dur_ms) in enumerate(queue):
        if i == 0:
            node_type = 'clickEffect'
            actual_delay = 0
        else:
            node_type = 'afterEffect'
            actual_delay = max(0, delay_ms - prev_end_ms)
        prev_end_ms = delay_ms + dur_ms
        anim_par = _build_entrance_anim(shape, anim_type, actual_delay, dur_ms, node_type)
        group_child.append(anim_par)

    prevCond = etree.SubElement(seq, qn('p:prevCondLst'))
    pc = etree.SubElement(prevCond, qn('p:cond'))
    pc.set('evt', 'onPrev')
    pc.set('delay', '0')
    pc_tgt = etree.SubElement(pc, qn('p:tgtEl'))
    etree.SubElement(pc_tgt, qn('p:sldTgt'))

    nextCond = etree.SubElement(seq, qn('p:nextCondLst'))
    nc = etree.SubElement(nextCond, qn('p:cond'))
    nc.set('evt', 'onNext')
    nc.set('delay', '0')
    nc_tgt = etree.SubElement(nc, qn('p:tgtEl'))
    etree.SubElement(nc_tgt, qn('p:sldTgt'))


# =====================================================
# SLIDES — v5 (25 pages)
# =====================================================

# ===== Slide 1: 封面 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_line = add_rect(s, 1.5, 2.2, 4, 0.04, fill_color=C['accent'])
shape_title = add_text(s, 1.5, 2.5, 10, 1.5, 'Kimi K3', size=54, bold=True, color=C['primary'], font=FONT_EN)
shape_sub = add_text(s, 1.5, 3.8, 10, 0.8, '当开源模型第一次逼到闭源喉咙口', size=22, color=C['text'])
shape_brand = add_text(s, 1.5, 5.0, 10, 0.5, 'AI大家学 · 前沿AI认知篇 · 2026年8月', size=14, color=C['subtext'])
shape_circle = add_ellipse(s, 9.5, 1.0, 3.5, 3.5, fill_color=C['card'], border_color=C['border'], border_width=1.5)
shape_num = add_text(s, 9.5, 2.0, 3.5, 0.8, '2.8T', size=36, bold=True, color=C['accent'], font=FONT_EN, align=PP_ALIGN.CENTER)
shape_label = add_text(s, 9.5, 2.8, 3.5, 0.5, '参数', size=14, color=C['subtext'], align=PP_ALIGN.CENTER)
add_footer(s, '数据来源：Artificial Analysis · Frontend Code Arena · 月之暗面官方技术报告')
animate(s, shape_line, 'wipe_right', 0.0, 0.5)
animate(s, shape_title, 'fade', 0.3, 0.8)
animate(s, shape_sub, 'fade', 0.8, 0.6)
animate(s, shape_brand, 'fade', 1.2, 0.5)
animate(s, shape_circle, 'zoom', 0.5, 0.8)
animate(s, shape_num, 'fade', 1.0, 0.6)
animate(s, shape_label, 'fade', 1.2, 0.5)
flush_animations(s)

# ===== Slide 2: Hook - 4700亿 + 融资 + HF =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_72h = add_text(s, 1.5, 1.0, 5, 1.0, '72小时', size=44, bold=True, color=C['primary'], font=FONT_EN)
shape_line2 = add_rect(s, 1.5, 1.9, 8, 0.03, fill_color=C['accent'])
shape_money = add_text(s, 1.5, 2.1, 10, 1.2, '-$4,700 亿', size=56, bold=True, color=C['red'], font=FONT_EN)
shape_evap = add_text(s, 1.5, 3.3, 10, 0.5, '全球科技股市值蒸发', size=18, color=C['text'])
shape_d1 = add_text_multi(s, 1.5, 4.1, 10, 0.4, [
    ('7月16日  ', 13, False, C['subtext'], FONT_CN),
    ('月之暗面发布 Kimi K3', 13, True, C['primary'], FONT_CN),
])
shape_d2 = add_text_multi(s, 1.5, 4.5, 10, 0.4, [
    ('7月22日  ', 13, False, C['subtext'], FONT_CN),
    ('白宫指控蒸馏 · 商务部调查 · 200家硅谷公司反对', 13, True, C['red'], FONT_CN),
])
shape_d3 = add_text_multi(s, 1.5, 4.9, 10, 0.4, [
    ('7月27日  ', 13, False, C['subtext'], FONT_CN),
    ('开源权重上线 · 30分钟登顶 Hugging Face', 13, True, C['teal'], FONT_CN),
])
shape_d4 = add_text_multi(s, 1.5, 5.3, 10, 0.4, [
    ('7月30日  ', 13, False, C['subtext'], FONT_CN),
    ('F轮融资 $35亿 · 估值飙至 $350亿 · 冲刺港股IPO', 13, True, C['accent'], FONT_CN),
])
add_footer(s, '据多家财经媒体统计 · Bloomberg · 36氪')
animate(s, shape_72h, 'wipe_right', 0.0, 0.6)
animate(s, shape_line2, 'wipe_right', 0.4, 0.5)
animate(s, shape_money, 'zoom', 0.8, 0.8)
animate(s, shape_evap, 'fade', 1.5, 0.5)
animate(s, shape_d1, 'fade', 2.0, 0.5)
animate(s, shape_d2, 'fade', 2.5, 0.5)
animate(s, shape_d3, 'fade', 3.0, 0.5)
animate(s, shape_d4, 'fade', 3.5, 0.5)
flush_animations(s)

# ===== Slide 3: 章节 - 它有多强 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
add_section_page(s, '01', '它有多强', '2.8万亿参数 · MoE 896专家 · 100万Token上下文')
flush_animations(s)

# ===== Slide 4: MoE 架构 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 8, 0.7, 'MoE 混合专家架构', size=28, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, '知识储备极大，单次推理只调用 1.8% 的专家', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

moe_data = [
    ('总参数', '2.8万亿', C['primary']),
    ('专家总数', '896', C['teal']),
    ('每次激活', '16个', C['accent']),
    ('激活比例', '1.8%', C['accent']),
]
for i, (label, value, color) in enumerate(moe_data):
    x = 0.8 + i * 3.1
    card = add_rect(s, x, 1.8, 2.8, 1.6, fill_color=C['card'], border_color=C['border'], radius=True)
    lb = add_text(s, x, 1.95, 2.8, 0.5, label, size=12, color=C['subtext'], align=PP_ALIGN.CENTER)
    vl = add_text(s, x, 2.4, 2.8, 0.8, value, size=28, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.CENTER)
    animate(s, card, 'wipe_bottom', 0.5 + i*0.4, 0.5)
    animate(s, lb, 'fade', 0.7 + i*0.4, 0.4)
    animate(s, vl, 'fade', 0.8 + i*0.4, 0.4)

explain_bg = add_rect(s, 0.8, 3.8, 11.73, 2.8, fill_color=C['card'], border_color=C['border'], radius=True)
explain_title = add_text(s, 1.2, 3.95, 3, 0.5, '翻译成人话', size=16, bold=True, color=C['accent'])
explain_body = add_text(s, 1.2, 4.5, 11, 2.0,
    '896个专家随时待命，但每次推理只叫16个出来干活。\n\n就像一家有896个专科医生的医院，你每次看病只需要挂16个相关科室。\n\n→ 知识储备极大，但单次成本很低',
    size=14, color=C['text'], line_spacing=20)
add_footer(s, '来源：月之暗面官方技术报告')
animate(s, explain_bg, 'wipe_bottom', 2.2, 0.6)
animate(s, explain_title, 'fade', 2.5, 0.4)
animate(s, explain_body, 'fade', 2.8, 0.8)
flush_animations(s)

# ===== Slide 5: 榜单成绩 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 8, 0.7, '成绩单 — 硬指标', size=28, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

benchmarks = [
    ('Frontend Code Arena', '前端代码能力', '1679', '#1', '全球第一 · 开源首次登顶', C['teal']),
    ('AA Intelligence Index', '智能指数', '57', '#3', '仅次 Fable 5 (60) 和 GPT-5.6 Sol (59)', C['accent']),
    ('BrowseComp', '长周期检索', '91.2', '#1', '全球第一', C['purple']),
    ('GDPval-AA v2', '真实经济任务', '1668', '#4', '开源模型最高', C['primary']),
]
for i, (name, sub, score, rank, note, color) in enumerate(benchmarks):
    y = 1.4 + i * 1.3
    bg = add_rect(s, 0.8, y, 11.73, 1.1, fill_color=C['card'], border_color=C['border'], radius=True)
    bar = add_rect(s, 0.8, y, 0.08, 1.1, fill_color=color)
    nm = add_text(s, 1.2, y+0.1, 4, 0.4, name, size=14, bold=True, color=C['text'], font=FONT_EN)
    sb = add_text(s, 1.2, y+0.5, 4, 0.4, sub, size=11, color=C['subtext'])
    sc = add_text(s, 7, y+0.15, 3, 0.7, score, size=28, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.RIGHT)
    badge = add_rect(s, 10.2, y+0.25, 0.9, 0.5, fill_color=color, radius=True)
    rk = add_text(s, 10.2, y+0.25, 0.9, 0.5, rank, size=14, bold=True, color=C['white'], font=FONT_EN, align=PP_ALIGN.CENTER)
    nt = add_text(s, 11.3, y+0.3, 2, 0.4, note, size=10, color=C['subtext'])
    animate(s, bg, 'wipe_right', 0.3 + i*0.6, 0.5)
    animate(s, bar, 'wipe_bottom', 0.5 + i*0.6, 0.3)
    animate(s, nm, 'fade', 0.6 + i*0.6, 0.4)
    animate(s, sb, 'fade', 0.7 + i*0.6, 0.4)
    animate(s, sc, 'fade', 0.8 + i*0.6, 0.4)
    animate(s, badge, 'zoom', 0.9 + i*0.6, 0.3)
    animate(s, rk, 'fade', 1.0 + i*0.6, 0.3)
    animate(s, nt, 'fade', 1.0 + i*0.6, 0.3)
add_footer(s, 'Artificial Analysis · Frontend Code Arena · 月之暗面技术报告')
flush_animations(s)

# ===== Slide 6: 三模型对比 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 8, 0.7, '三模型对比', size=28, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

col_x = [0.8, 4.8, 8.0, 11.0]
col_w = [3.8, 3.0, 2.8, 1.73]
models = [
    ('对比维度', 'Code Arena', 'AA Index', '开源/闭源', C['subtext'], True),
    ('Kimi K3', '1679', '57', '开源', C['teal'], False),
    ('Claude Fable 5', '1631', '60', '闭源', C['accent'], False),
    ('GPT-5.6 Sol', '1618', '59', '闭源', C['red'], False),
]
for i, (name, ca, aa, oc, color, is_header) in enumerate(models):
    y = 1.4 + i * 1.2
    bg_c = C['primary'] if is_header else C['card']
    tc = C['white'] if is_header else C['text']
    bg = add_rect(s, 0.8, y, 11.73, 1.0, fill_color=bg_c,
                  border_color=C['border'] if not is_header else None, radius=True)
    n = add_text(s, col_x[0], y+0.15, col_w[0], 0.7, name,
                 size=13 if is_header else 16, bold=not is_header, color=C['white'] if is_header else color)
    c1 = add_text(s, col_x[1], y+0.15, col_w[1], 0.7, ca,
                  size=12 if is_header else 18, bold=not is_header, color=tc, align=PP_ALIGN.CENTER, font=FONT_EN)
    c2 = add_text(s, col_x[2], y+0.15, col_w[2], 0.7, aa,
                  size=12 if is_header else 18, bold=not is_header, color=tc, align=PP_ALIGN.CENTER, font=FONT_EN)
    c3 = add_text(s, col_x[3], y+0.15, col_w[3], 0.7, oc,
                  size=12 if is_header else 13, color=C['white'] if is_header else color, align=PP_ALIGN.CENTER)
    animate(s, bg, 'wipe_right', 0.2 + i*0.5, 0.5)
    animate(s, n, 'fade', 0.4 + i*0.5, 0.4)
    animate(s, c1, 'fade', 0.5 + i*0.5, 0.4)
    animate(s, c2, 'fade', 0.6 + i*0.5, 0.4)
    animate(s, c3, 'fade', 0.7 + i*0.5, 0.4)

conc_bg = add_rect(s, 0.8, 6.0, 11.73, 0.7, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
conc = add_text(s, 0.8, 6.05, 11.73, 0.6, '开源第一，综合第三，性价比维度几乎没有对手',
                size=16, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)
animate(s, conc_bg, 'wipe_bottom', 2.5, 0.5)
animate(s, conc, 'fade', 2.8, 0.5)
flush_animations(s)

# ===== Slide 7: I-Love-Q 时钟对比 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '真实生产力 — I-Love-Q 中子星研究', size=24, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

k3_bg = add_rect(s, 0.8, 1.4, 5.5, 5.0, fill_color=C['card'], border_color=C['teal'], border_width=2, radius=True)
k3_hdr = add_rect(s, 0.8, 1.4, 5.5, 0.6, fill_color=C['teal'], radius=True)
k3_name = add_text(s, 0.8, 1.42, 5.5, 0.55, 'Kimi K3', size=16, bold=True, color=C['white'], font=FONT_EN, align=PP_ALIGN.CENTER)
k3_time = add_text(s, 0.8, 2.5, 5.5, 1.0, '~2 小时', size=44, bold=True, color=C['teal'], font=FONT_EN, align=PP_ALIGN.CENTER)
k3_desc = add_text(s, 1.2, 3.6, 4.7, 2.5,
    '自主阅读 20+ 篇论文\n评估 300+ 种状态方程\n生成 3000+ 行 Python\n马斯克评价："Impressive"',
    size=13, color=C['text'], line_spacing=22)

vs = add_text(s, 5.8, 3.5, 1.5, 0.8, 'VS', size=20, bold=True, color=C['subtext'], font=FONT_EN, align=PP_ALIGN.CENTER)

h_bg = add_rect(s, 7, 1.4, 5.5, 5.0, fill_color=C['card'], border_color=C['accent'], border_width=2, radius=True)
h_hdr = add_rect(s, 7, 1.4, 5.5, 0.6, fill_color=C['accent'], radius=True)
h_name = add_text(s, 7, 1.42, 5.5, 0.55, '资深研究人员', size=16, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
h_time = add_text(s, 7, 2.5, 5.5, 1.0, '1-2 周', size=44, bold=True, color=C['accent'], font=FONT_EN, align=PP_ALIGN.CENTER)
h_desc = add_text(s, 7.4, 3.6, 4.7, 2.5,
    '需要领域专业知识积累\n手动检索和筛选文献\n编写和调试计算代码\n反复验证状态方程',
    size=13, color=C['text'], line_spacing=22)

add_footer(s, '来源：月之暗面官方演示')
animate(s, k3_bg, 'fly_left', 0.3, 0.7)
animate(s, k3_hdr, 'wipe_right', 0.8, 0.3)
animate(s, k3_name, 'fade', 1.0, 0.4)
animate(s, k3_time, 'zoom', 1.2, 0.6)
animate(s, k3_desc, 'fade', 1.6, 0.8)
animate(s, vs, 'zoom', 2.0, 0.5)
animate(s, h_bg, 'fly_right', 2.3, 0.7)
animate(s, h_hdr, 'wipe_right', 2.8, 0.3)
animate(s, h_name, 'fade', 3.0, 0.4)
animate(s, h_time, 'zoom', 3.2, 0.6)
animate(s, h_desc, 'fade', 3.6, 0.8)
flush_animations(s)

# ===== Slide 8: 章节 - 怎么做到的 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
add_section_page(s, '02', '怎么做到的', '不是堆参数，是改地基')
flush_animations(s)

# ===== Slide 9: 三大技术 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '三项自研底层技术', size=28, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, '对沿用近十年的 AI 基础机制做系统性重写', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

techs = [
    ('KDA 混合线性注意力', '不同通道 × 不同衰减速率\n管理长序列记忆\n+ Attention Residuals', '推理效率 ↑ ~25%', C['teal']),
    ('Kimi Linear Tension', '自研线性注意力机制\n解决超长任务性能衰减\n上下文窗口 ×10', '长上下文的工程底气', C['accent']),
    ('Moon Clip 二阶优化器', '首次用于大模型训练\n20T 数据跑出 40T 效果\n训练成本和算力功耗减半', '训练效率翻倍', C['purple']),
]
for i, (name, desc, result, color) in enumerate(techs):
    x = 0.8 + i * 4.1
    card = add_rect(s, x, 1.7, 3.8, 4.5, fill_color=C['card'], border_color=C['border'], radius=True)
    hdr = add_rect(s, x, 1.7, 3.8, 0.55, fill_color=color, radius=True)
    nm = add_text(s, x, 1.72, 3.8, 0.5, name, size=14, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    ds = add_text(s, x+0.3, 2.6, 3.2, 2.5, desc, size=12, color=C['text'], line_spacing=20)
    sep = add_rect(s, x+0.3, 4.8, 3.2, 0.02, fill_color=C['border'])
    rs = add_text(s, x, 5.0, 3.8, 0.6, result, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    animate(s, card, 'wipe_bottom', 0.5 + i*0.7, 0.6)
    animate(s, hdr, 'wipe_right', 0.8 + i*0.7, 0.3)
    animate(s, nm, 'fade', 0.9 + i*0.7, 0.4)
    animate(s, ds, 'fade', 1.1 + i*0.7, 0.5)
    animate(s, sep, 'wipe_right', 1.3 + i*0.7, 0.3)
    animate(s, rs, 'fade', 1.4 + i*0.7, 0.4)

res_bg = add_rect(s, 2, 6.4, 9.33, 0.6, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
res = add_text(s, 2, 6.42, 9.33, 0.55, '三项叠加 → 整体扩展效率提升 ~2.5 倍',
               size=16, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)
animate(s, res_bg, 'wipe_bottom', 3.8, 0.5)
animate(s, res, 'fade', 4.1, 0.5)
flush_animations(s)

# ===== Slide 10: 章节 - 为什么全球紧张 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
add_section_page(s, '03', '为什么全球都紧张', '算力挤兑 · 蒸馏争议 · 从卡芯片到卡算法')
flush_animations(s)

# ===== Slide 11: 蒸馏争议 (v5 更新) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '蒸馏争议 — 三条反驳', size=28, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

acc_bg = add_rect(s, 0.8, 1.3, 11.73, 1.1, fill_color=RGBColor(0xFD,0xF0,0xF0), border_color=C['red'], radius=True)
acc = add_text(s, 1.2, 1.35, 11, 1.0,
    '白宫指控：月之暗面蒸馏了 Anthropic 的 Fable 模型\n财长威胁制裁 · 商务部 BIS 确认调查 · 截至8月初尚未列入实体清单',
    size=13, color=C['red'], line_spacing=19)
animate(s, acc_bg, 'wipe_right', 0.3, 0.5)
animate(s, acc, 'fade', 0.6, 0.5)

refutations = [
    ('1', '性能无法靠蒸馏获得', 'K3 排全球第二/第三/前端第一，仅靠蒸馏不可能', C['teal']),
    ('2', '开源技术报告可复核', '开放权重 + 全公开架构，蒸馏出来的是别人的架构', C['accent']),
    ('3', '时间线对不上', 'Fable 发布到 K3 训练完成之间不够走完蒸馏流程', C['purple']),
]
for i, (num, title, desc, color) in enumerate(refutations):
    y = 2.7 + i * 1.2
    bg = add_rect(s, 0.8, y, 11.73, 1.0, fill_color=C['card'], border_color=C['border'], radius=True)
    circle = add_ellipse(s, 1.1, y+0.15, 0.6, 0.6, fill_color=color)
    nm = add_text(s, 1.1, y+0.15, 0.6, 0.6, num, size=18, bold=True, color=C['white'], font=FONT_EN, align=PP_ALIGN.CENTER)
    tl = add_text(s, 2.1, y+0.1, 6, 0.4, title, size=15, bold=True, color=C['text'])
    ds = add_text(s, 2.1, y+0.5, 9, 0.4, desc, size=11, color=C['subtext'])
    chk = add_text(s, 11.5, y+0.2, 0.8, 0.6, '✓', size=22, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.CENTER)
    animate(s, bg, 'wipe_right', 1.2 + i*0.7, 0.5)
    animate(s, circle, 'zoom', 1.5 + i*0.7, 0.3)
    animate(s, nm, 'fade', 1.6 + i*0.7, 0.3)
    animate(s, tl, 'fade', 1.7 + i*0.7, 0.4)
    animate(s, ds, 'fade', 1.8 + i*0.7, 0.4)
    animate(s, chk, 'zoom', 2.0 + i*0.7, 0.3)

# v5: 200家硅谷反对 + 中方反控
extra_bg = add_rect(s, 0.8, 6.3, 11.73, 0.6, fill_color=C['card'], border_color=C['teal'], border_width=1.5, radius=True)
extra = add_text(s, 1.2, 6.33, 11, 0.55,
    '反转：近200家硅谷公司联名反对制裁 · 中国商务部反控美方也在蒸馏中国模型',
    size=13, bold=True, color=C['teal'], align=PP_ALIGN.CENTER)
animate(s, extra_bg, 'wipe_bottom', 4.5, 0.5)
animate(s, extra, 'fade', 4.8, 0.5)
add_footer(s, '来源：Nathan Lambert · BBC中文 · 电子工程专辑')
flush_animations(s)

# ===== Slide 12: 卡芯片到卡算法 (v5 更新 — 5节点) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '从"卡芯片"到"卡算法"', size=28, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

axis = add_rect(s, 0.6, 3.0, 12.13, 0.03, fill_color=C['border'])
animate(s, axis, 'wipe_right', 0.3, 0.6)

phases = [
    ('禁运GPU', '2022-24', '切断高端\n芯片供应', C['red']),
    ('指控蒸馏', '2025-26', 'DeepSeek\n+ Kimi K3', C['accent']),
    ('商务部调查', '2026.07', 'BIS确认\n调查GB300', C['red']),
    ('硅谷反对', '2026.07', '200家公司\n联名反对', C['teal']),
    ('中方反控', '2026.07', '美方也在\n蒸馏中国模型', C['purple']),
]
for i, (name, year, desc, color) in enumerate(phases):
    x = 0.6 + i * 2.5
    card = add_rect(s, x, 1.5, 2.2, 1.2, fill_color=C['card'], border_color=C['border'], radius=True)
    nm = add_text(s, x, 1.55, 2.2, 0.45, name, size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    yr = add_text(s, x, 1.95, 2.2, 0.3, year, size=9, color=C['subtext'], font=FONT_EN, align=PP_ALIGN.CENTER)
    dot = add_ellipse(s, x+0.5, 2.85, 0.3, 0.3, fill_color=color)
    ds = add_text(s, x, 3.3, 2.2, 0.6, desc, size=10, color=C['subtext'], align=PP_ALIGN.CENTER, line_spacing=14)
    animate(s, card, 'wipe_bottom', 0.6 + i*0.4, 0.5)
    animate(s, nm, 'fade', 0.8 + i*0.4, 0.4)
    animate(s, yr, 'fade', 0.9 + i*0.4, 0.3)
    animate(s, dot, 'zoom', 1.0 + i*0.4, 0.3)
    animate(s, ds, 'fade', 1.1 + i*0.4, 0.3)

analysis_bg = add_rect(s, 0.8, 4.3, 11.73, 2.5, fill_color=C['card'], border_color=C['border'], radius=True)
analysis_t = add_text(s, 1.2, 4.4, 11, 0.5, 'Pattern：你强了，就说你偷', size=16, bold=True, color=C['primary'])
analysis_b = add_text(s, 1.2, 4.95, 11, 1.7,
    '• 6月：Anthropic 指控阿里巴巴"蒸馏攻击"\n• 7月：白宫指控月之暗面蒸馏 Fable · BIS 展开调查\n• 同日：近200家硅谷公司联名反对制裁中国AI模型\n• 中方反控：美国AI企业也在蒸馏中国模型\n\n连硅谷自己人都不买账 — 制裁开源只会把全球开发者推向中国生态',
    size=12, color=C['text'], line_spacing=18)
animate(s, analysis_bg, 'wipe_bottom', 3.5, 0.6)
animate(s, analysis_t, 'fade', 3.8, 0.5)
animate(s, analysis_b, 'fade', 4.1, 0.8)
flush_animations(s)

# ===== Slide 13: 章节 - 为什么重要 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
add_section_page(s, '04', '为什么重要', '趋势信号 · 开源突围 · 商业数据 · 资本定价')
flush_animations(s)

# ===== Slide 14: 趋势时间轴 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '中国AI突破密度 — 不是偶然，是趋势', size=24, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

axis = add_rect(s, 0.8, 3.3, 11.73, 0.03, fill_color=C['border'])
animate(s, axis, 'wipe_right', 0.3, 0.8)

events = [
    ('DeepSeek R1', '2025.01', '推理模型\n震惊华尔街', C['accent'], False),
    ('DeepSeek V3', '2025年中', '开源基座\n逼近GPT-4', C['accent'], False),
    ('Kimi K2', '2025年底', 'MoE架构\n1万亿参数', C['teal'], False),
    ('Kimi K3', '2026.07', '2.8万亿\n全球第三', C['teal'], True),
]
for i, (name, date, desc, color, highlight) in enumerate(events):
    x = 1.2 + i * 3.1
    y_top = 1.2 if highlight else 1.5
    h_card = 1.8 if highlight else 1.5
    bw = 2 if highlight else 1
    bc = color if highlight else C['border']
    card = add_rect(s, x, y_top, 2.8, h_card, fill_color=C['card'], border_color=bc, border_width=bw, radius=True)
    nm = add_text(s, x, y_top+0.05, 2.8, 0.45, name, size=16 if highlight else 14, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.CENTER)
    dt = add_text(s, x, y_top+0.45, 2.8, 0.3, date, size=10, color=C['subtext'], font=FONT_EN, align=PP_ALIGN.CENTER)
    ds = add_text(s, x, y_top+0.8, 2.8, 0.7, desc, size=11, color=C['text'], align=PP_ALIGN.CENTER, line_spacing=16)
    dot = add_ellipse(s, x+0.8, 3.15, 0.3, 0.3, fill_color=color)
    animate(s, card, 'wipe_bottom', 0.6 + i*0.6, 0.5)
    animate(s, nm, 'fade', 0.8 + i*0.6, 0.4)
    animate(s, dt, 'fade', 0.9 + i*0.6, 0.3)
    animate(s, ds, 'fade', 1.0 + i*0.6, 0.4)
    animate(s, dot, 'zoom', 1.1 + i*0.6, 0.3)

conc_bg = add_rect(s, 1.5, 4.0, 10.33, 2.5, fill_color=C['card'], border_color=C['border'], radius=True)
conc_t = add_text(s, 6.66, 4.1, 4, 0.5, '关键判断', size=16, bold=True, color=C['primary'], align=PP_ALIGN.CENTER)
conc_b = add_text(s, 2, 4.65, 9.33, 1.7,
    'DeepSeek 出来时，有人说是灵光一现。K3 出来，你不能再这么说了。\n\n中国AI公司正在系统性逼近、甚至局部超越美国顶级实验室。\n不是一次，是一种趋势。',
    size=13, color=C['text'], align=PP_ALIGN.CENTER, line_spacing=20)
animate(s, conc_bg, 'wipe_bottom', 3.2, 0.6)
animate(s, conc_t, 'fade', 3.5, 0.4)
animate(s, conc_b, 'fade', 3.8, 0.8)
flush_animations(s)

# ===== Slide 15: 开源三件套 (v5 新增) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '开源不只是模型', size=28, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, '给了菜谱，还给了整套厨房设备 — 模型 + 训练基础设施全链条开源', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

infra = [
    ('MoonEP', '高性能通信库', '解决大规模 MoE 模型\n训练中的分布式\n通信瓶颈', '通信', C['teal']),
    ('FlashKDA', '高性能注意力算子', '为 KDA 提供计算内核\nPrefill 提速\n1.7 - 2.2 倍', '算子', C['accent']),
    ('AgentEnv', 'Agent训练沙箱', '分布式强化学习环境\n降低大规模 Agent\nRL 训练开销', '沙箱', C['purple']),
]
for i, (name, sub, desc, tag, color) in enumerate(infra):
    x = 0.8 + i * 4.1
    card = add_rect(s, x, 1.7, 3.8, 3.5, fill_color=C['card'], border_color=C['border'], radius=True)
    hdr = add_rect(s, x, 1.7, 3.8, 0.55, fill_color=color, radius=True)
    nm = add_text(s, x, 1.72, 3.8, 0.5, name, size=16, bold=True, color=C['white'], font=FONT_EN, align=PP_ALIGN.CENTER)
    sb = add_text(s, x, 2.5, 3.8, 0.4, sub, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    ds = add_text(s, x+0.3, 3.0, 3.2, 1.8, desc, size=12, color=C['text'], align=PP_ALIGN.CENTER, line_spacing=18)
    tag_bg = add_rect(s, x+1.2, 4.5, 1.4, 0.35, fill_color=color, radius=True)
    tag_t = add_text(s, x+1.2, 4.5, 1.4, 0.35, tag, size=10, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    animate(s, card, 'wipe_bottom', 0.5 + i*0.6, 0.6)
    animate(s, hdr, 'wipe_right', 0.8 + i*0.6, 0.3)
    animate(s, nm, 'fade', 0.9 + i*0.6, 0.4)
    animate(s, sb, 'fade', 1.0 + i*0.6, 0.3)
    animate(s, ds, 'fade', 1.1 + i*0.6, 0.5)
    animate(s, tag_bg, 'zoom', 1.3 + i*0.6, 0.3)
    animate(s, tag_t, 'fade', 1.4 + i*0.6, 0.3)

# 底部数据条
stats_bg = add_rect(s, 0.8, 5.6, 11.73, 1.1, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
stat_items = [
    ('30分钟', '登顶HF趋势榜', 1.3, C['teal']),
    ('113万+', 'HuggingFace下载', 4.3, C['accent']),
    ('MIT许可', '免费商用', 7.3, C['purple']),
    ('Day 0', '华为/阿里/vLLM适配', 10.0, C['primary']),
]
for val, label, x, color in stat_items:
    vl = add_text(s, x, 5.65, 2.2, 0.5, val, size=18, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.CENTER)
    lb = add_text(s, x, 6.1, 2.2, 0.4, label, size=10, color=C['subtext'], align=PP_ALIGN.CENTER)
    animate(s, vl, 'fade', 2.5, 0.4)
    animate(s, lb, 'fade', 2.7, 0.4)
animate(s, stats_bg, 'wipe_bottom', 2.3, 0.5)
add_footer(s, '来源：Hugging Face · vLLM · 华为昇腾 · 财联社')
flush_animations(s)

# ===== Slide 16: 生态闭环 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '中国AI基础设施闭环', size=28, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

nodes = [
    ('国产芯片', '华为昇腾 · 沐曦 · 燧原', 4.67, 1.5, C['red'], 'fly_top'),
    ('开源模型', 'Kimi K3 · DeepSeek', 8.5, 3.5, C['teal'], 'fly_right'),
    ('部署框架', 'vLLM · SGLang · GPUStack', 4.67, 5.5, C['accent'], 'fly_bottom'),
    ('应用生态', '企业应用 · 开发者', 0.8, 3.5, C['purple'], 'fly_left'),
]
for i, (name, sub, x, y, color, fly_dir) in enumerate(nodes):
    card = add_rect(s, x, y, 4, 1.3, fill_color=C['card'], border_color=color, border_width=1.5, radius=True)
    nm = add_text(s, x, y+0.15, 4, 0.5, name, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    sb = add_text(s, x, y+0.65, 4, 0.4, sub, size=11, color=C['subtext'], align=PP_ALIGN.CENTER)
    animate(s, card, fly_dir, 0.3 + i*0.5, 0.5)
    animate(s, nm, 'fade', 0.6 + i*0.5, 0.4)
    animate(s, sb, 'fade', 0.7 + i*0.5, 0.4)

center = add_text(s, 5.17, 3.3, 3, 1.2, '自主可控\n完整闭环', size=18, bold=True, color=C['primary'], align=PP_ALIGN.CENTER, line_spacing=26)
animate(s, center, 'zoom', 2.5, 0.6)
add_footer(s, '华为昇腾、阿里云已Day 0适配 K3 · vLLM 提供生产级支持')
flush_animations(s)

# ===== Slide 17: 商业数据爆发 (v5 新增) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, 'K3 发布后的商业数据', size=28, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, '日销售额环比增长至少6倍 · K3上线次日ARR创历史最大单日增幅', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

biz_data = [
    ('×6', '日销售额增长', 'K3发布后环比暴增', C['red']),
    ('$3亿+', '年化收入(ARR)', '6月中旬已突破\nAPI占比超70%', C['teal']),
    ('70%+', 'API收入占比', '海外收入已超国内\nB端规模化变现', C['accent']),
]
for i, (num, title, desc, color) in enumerate(biz_data):
    x = 0.8 + i * 4.1
    card = add_rect(s, x, 1.7, 3.8, 3.0, fill_color=C['card'], border_color=color, border_width=2, radius=True)
    big = add_text(s, x, 1.9, 3.8, 0.9, num, size=40, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.CENTER)
    tl = add_text(s, x, 2.85, 3.8, 0.4, title, size=14, bold=True, color=C['text'], align=PP_ALIGN.CENTER)
    ds = add_text(s, x+0.3, 3.3, 3.2, 1.2, desc, size=11, color=C['subtext'], align=PP_ALIGN.CENTER, line_spacing=16)
    animate(s, card, 'wipe_bottom', 0.5 + i*0.6, 0.6)
    animate(s, big, 'zoom', 0.9 + i*0.6, 0.5)
    animate(s, tl, 'fade', 1.2 + i*0.6, 0.4)
    animate(s, ds, 'fade', 1.4 + i*0.6, 0.5)

# 融资条
fund_bg = add_rect(s, 0.8, 5.2, 11.73, 1.4, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
fund_t = add_text(s, 1.2, 5.3, 11, 0.5, 'F轮超额融资', size=16, bold=True, color=C['accent'])
fund_b = add_text(s, 1.2, 5.85, 11, 0.65,
    '原定目标 $10-20亿 → 实际超募 $35亿（近3倍）· 提前关闭\n投后估值从 $315亿 跳至 $350亿 · 拆除VIE · 冲刺港股IPO',
    size=12, color=C['text'], line_spacing=18)
animate(s, fund_bg, 'wipe_bottom', 2.8, 0.5)
animate(s, fund_t, 'fade', 3.1, 0.4)
animate(s, fund_b, 'fade', 3.3, 0.5)
add_footer(s, '来源：IT之家（张予彤发言）· 36氪 · 财联社 · 中国风投网')
flush_animations(s)

# ===== Slide 18: 估值 (v5 更新 — 加350亿F轮) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '估值飙升 — 真金白银的市场定价', size=24, bold=True, color=C['primary'])
animate(s, shape_t, 'fade', 0, 0.5)

axis = add_rect(s, 0.6, 3.5, 12.13, 0.03, fill_color=C['border'])
animate(s, axis, 'wipe_right', 0.3, 0.6)

vals = [
    ('2024.08', '$33亿', 0.6, False),
    ('2025.12', '$43亿', 2.8, False),
    ('2026.01', '$100亿+', 5.0, False),
    ('2026.05', '$200亿', 7.2, False),
    ('2026.06', '$315亿', 9.2, False),
    ('2026.07', '$350亿', 11.2, True),
]
for i, (date, val, x, hl) in enumerate(vals):
    c = C['accent'] if hl else C['primary']
    sz = 20 if hl else 14
    dot = add_ellipse(s, x+0.35, 3.35, 0.3, 0.3, fill_color=c)
    vl = add_text(s, x, 2.5, 1.5, 0.6, val, size=sz, bold=True, color=c, font=FONT_EN, align=PP_ALIGN.CENTER)
    dt = add_text(s, x, 3.8, 1.5, 0.4, date, size=9, color=C['subtext'], font=FONT_EN, align=PP_ALIGN.CENTER)
    animate(s, dot, 'zoom', 0.6 + i*0.35, 0.3)
    animate(s, vl, 'fade', 0.7 + i*0.35, 0.4)
    animate(s, dt, 'fade', 0.8 + i*0.35, 0.3)

# Pre-IPO target
target_bg = add_rect(s, 10.5, 1.3, 2.5, 1.5, fill_color=C['card'], border_color=C['purple'], border_width=1.5, radius=True)
target_v = add_text(s, 10.5, 1.4, 2.5, 0.6, '$500亿?', size=20, bold=True, color=C['purple'], font=FONT_EN, align=PP_ALIGN.CENTER)
target_d = add_text(s, 10.5, 2.0, 2.5, 0.7, 'Pre-IPO目标\n中金+高盛 · 港股\n拆除VIE架构', size=10, color=C['subtext'], align=PP_ALIGN.CENTER, line_spacing=14)
animate(s, target_bg, 'wipe_right', 3.0, 0.5)
animate(s, target_v, 'fade', 3.3, 0.4)
animate(s, target_d, 'fade', 3.5, 0.4)

conc_bg = add_rect(s, 2, 4.8, 9.33, 1.8, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
conc = add_text(s, 2, 4.9, 9.33, 1.6,
    '从 $33亿 到 $350亿 · 半年涨约八倍\nF轮超募3倍 · 冲刺港股IPO\n"不该被贴上低价标签" — 从便宜到敢定价',
    size=16, bold=True, color=C['accent'], align=PP_ALIGN.CENTER, line_spacing=24)
animate(s, conc_bg, 'wipe_bottom', 3.8, 0.5)
animate(s, conc, 'fade', 4.1, 0.5)
add_footer(s, 'Bloomberg · 36氪 · 财联社 · 电子工程专辑')
flush_animations(s)

# ===== Slide 19: 章节 - 本地部署 (v5 新增) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
add_section_page(s, '05', '本地部署到底要多少钱？', '64卡超节点 · 8×B300 · 显存分析 · 成本评估')
flush_animations(s)

# ===== Slide 20: 64卡 vs 8卡 (v5 新增) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '本地部署 — 官方推荐 vs 单机最优', size=24, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, '模型权重 1560 GB · vLLM推理需要 ~1680 GB 显存', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

# 左侧：64卡
left_bg = add_rect(s, 0.8, 1.7, 5.5, 4.8, fill_color=C['card'], border_color=C['primary'], border_width=2, radius=True)
left_hdr = add_rect(s, 0.8, 1.7, 5.5, 0.55, fill_color=C['primary'], radius=True)
left_name = add_text(s, 0.8, 1.72, 5.5, 0.5, '64卡超节点（官方推荐）', size=14, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
left_gpu = add_text(s, 0.8, 2.6, 5.5, 0.6, '8台 × 8卡 H200', size=20, bold=True, color=C['primary'], font=FONT_EN, align=PP_ALIGN.CENTER)
left_vram = add_text(s, 0.8, 3.3, 5.5, 0.5, '聚合显存 5,120 GB', size=16, color=C['text'], font=FONT_EN, align=PP_ALIGN.CENTER)
left_detail = add_text(s, 1.2, 4.0, 4.7, 1.8,
    '✓ 充裕的KV缓存空间\n✓ 支持长上下文并发\n✓ 生产级稳定服务\n\n采购价 ~3000万人民币\n+网络运维 ≈ 4000万+',
    size=12, color=C['text'], line_spacing=18)

# VS
vs_txt = add_text(s, 5.8, 3.5, 1.5, 0.8, 'VS', size=20, bold=True, color=C['subtext'], font=FONT_EN, align=PP_ALIGN.CENTER)

# 右侧：8×B300
right_bg = add_rect(s, 7, 1.7, 5.5, 4.8, fill_color=C['card'], border_color=C['teal'], border_width=2, radius=True)
right_hdr = add_rect(s, 7, 1.7, 5.5, 0.55, fill_color=C['teal'], radius=True)
right_name = add_text(s, 7, 1.72, 5.5, 0.5, '8×B300 单机（最优解）', size=14, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
right_gpu = add_text(s, 7, 2.6, 5.5, 0.6, '8 × 288 GB HBM3e', size=20, bold=True, color=C['teal'], font=FONT_EN, align=PP_ALIGN.CENTER)
right_vram = add_text(s, 7, 3.3, 5.5, 0.5, '合计显存 2,304 GB', size=16, color=C['text'], font=FONT_EN, align=PP_ALIGN.CENTER)
right_detail = add_text(s, 7.4, 4.0, 4.7, 1.8,
    '✓ 1560GB权重 + 740GB缓存\n✓ 同一NVSwitch域（零跨节点）\n✓ FP4/FP8双精度硬件加速\n\n美国市场 ~300-400万人民币\n中国市场 ~700万+（出口管制溢价）',
    size=12, color=C['text'], line_spacing=18)

animate(s, left_bg, 'fly_left', 0.5, 0.7)
animate(s, left_hdr, 'wipe_right', 1.0, 0.3)
animate(s, left_name, 'fade', 1.2, 0.4)
animate(s, left_gpu, 'fade', 1.4, 0.4)
animate(s, left_vram, 'fade', 1.6, 0.3)
animate(s, left_detail, 'fade', 1.8, 0.6)
animate(s, vs_txt, 'zoom', 2.2, 0.5)
animate(s, right_bg, 'fly_right', 2.5, 0.7)
animate(s, right_hdr, 'wipe_right', 3.0, 0.3)
animate(s, right_name, 'fade', 3.2, 0.4)
animate(s, right_gpu, 'fade', 3.4, 0.4)
animate(s, right_vram, 'fade', 3.6, 0.3)
animate(s, right_detail, 'fade', 3.8, 0.6)
add_footer(s, '来源：vLLM · NVIDIA DGX B300 · Reuters · SegmentFault')
flush_animations(s)

# ===== Slide 21: B300 显存分配 (v5 新增) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '8×B300 显存分配', size=28, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, 'Blackwell Ultra 架构 · 单卡 288GB HBM3e · NVLink 5.0 (1800 GB/s)', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

# 显存分配条形图
total_bar = add_rect(s, 1.5, 2.0, 10.33, 1.5, fill_color=C['dark_card'], border_color=C['border'], radius=True)
total_label = add_text(s, 1.5, 2.05, 10.33, 0.4, '8×B300 总显存：2,304 GB', size=14, bold=True, color=C['text'], font=FONT_EN, align=PP_ALIGN.CENTER)
# 权重部分 (1560/2304 ≈ 67.7%)
w_width = 7.0  # ~67.7% of 10.33
w_bar = add_rect(s, 1.5, 2.5, w_width, 0.9, fill_color=C['teal'], radius=True)
w_label = add_text(s, 1.5, 2.6, w_width, 0.7, '模型权重 1,560 GB (MXFP4)', size=13, bold=True, color=C['white'], font=FONT_EN, align=PP_ALIGN.CENTER)
# 缓存部分 (744/2304 ≈ 32.3%)
c_width = 10.33 - w_width
c_bar = add_rect(s, 1.5 + w_width, 2.5, c_width, 0.9, fill_color=C['accent'], radius=True)
c_label = add_text(s, 1.5 + w_width, 2.55, c_width, 0.8, 'KV缓存\n+ 开销', size=11, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
animate(s, total_bar, 'wipe_right', 0.5, 0.6)
animate(s, total_label, 'fade', 0.8, 0.4)
animate(s, w_bar, 'wipe_right', 1.2, 0.6)
animate(s, w_label, 'fade', 1.5, 0.4)
animate(s, c_bar, 'wipe_right', 1.8, 0.5)
animate(s, c_label, 'fade', 2.1, 0.4)

# 详细数字卡片
detail_items = [
    ('288 GB', '单卡 HBM3e', C['teal']),
    ('2,304 GB', '8卡合计显存', C['primary']),
    ('1,560 GB', 'MXFP4 权重', C['teal']),
    ('~740 GB', 'KV缓存+开销', C['accent']),
]
for i, (val, label, color) in enumerate(detail_items):
    x = 0.8 + i * 3.1
    card = add_rect(s, x, 3.8, 2.8, 1.3, fill_color=C['card'], border_color=C['border'], radius=True)
    vl = add_text(s, x, 3.9, 2.8, 0.5, val, size=22, bold=True, color=color, font=FONT_EN, align=PP_ALIGN.CENTER)
    lb = add_text(s, x, 4.4, 2.8, 0.4, label, size=11, color=C['subtext'], align=PP_ALIGN.CENTER)
    animate(s, card, 'wipe_bottom', 2.5 + i*0.4, 0.5)
    animate(s, vl, 'fade', 2.8 + i*0.4, 0.4)
    animate(s, lb, 'fade', 2.9 + i*0.4, 0.3)

# 关键优势
adv_bg = add_rect(s, 0.8, 5.5, 11.73, 1.2, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
adv = add_text(s, 1.2, 5.55, 11, 1.1,
    'B300 核心优势：FP4/FP8 双精度硬件加速 · NVLink 5.0 带宽 1800 GB/s\n8卡全部在同一 NVSwitch 域内 → 避免跨节点通信开销 → 单机部署最优解',
    size=13, bold=True, color=C['accent'], align=PP_ALIGN.CENTER, line_spacing=20)
animate(s, adv_bg, 'wipe_bottom', 4.5, 0.5)
animate(s, adv, 'fade', 4.8, 0.5)
add_footer(s, '来源：NVIDIA DGX B300 · vLLM Blog · GPUStack 实测')
flush_animations(s)

# ===== Slide 22: 三条部署路线 (v5 新增) =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_t = add_text(s, 0.8, 0.4, 10, 0.7, '部署方案正在多元化', size=28, bold=True, color=C['primary'])
shape_st = add_text(s, 0.8, 1.0, 10, 0.4, '开源的意义不只是"免费用"，而是让部署选择权回到用户手里', size=13, color=C['subtext'])
animate(s, shape_t, 'fade', 0, 0.5)
animate(s, shape_st, 'fade', 0.3, 0.5)

routes = [
    ('NVIDIA B300', 'Blackwell Ultra', '单机8卡最优性能\nFP4/FP8硬件加速\nNVLink 1800GB/s\n\n~300-700万', '最强性能', C['teal']),
    ('AMD MI355X', 'Instinct', '单机8卡可跑K3\n性价比反超B300\nROCm生态持续追赶\n\n成本更低', '性价比之选', C['accent']),
    ('华为昇腾', 'Ascend 950', 'Day 0 完成适配\nCANN 全系列支持\n国产替代方案\n\n自主可控', '国产替代', C['purple']),
]
for i, (name, sub, desc, tag, color) in enumerate(routes):
    x = 0.8 + i * 4.1
    card = add_rect(s, x, 1.7, 3.8, 4.3, fill_color=C['card'], border_color=C['border'], radius=True)
    hdr = add_rect(s, x, 1.7, 3.8, 0.55, fill_color=color, radius=True)
    nm = add_text(s, x, 1.72, 3.8, 0.5, name, size=15, bold=True, color=C['white'], font=FONT_EN, align=PP_ALIGN.CENTER)
    sb = add_text(s, x, 2.5, 3.8, 0.35, sub, size=11, color=C['subtext'], font=FONT_EN, align=PP_ALIGN.CENTER)
    ds = add_text(s, x+0.3, 3.0, 3.2, 2.3, desc, size=12, color=C['text'], align=PP_ALIGN.CENTER, line_spacing=18)
    tag_bg = add_rect(s, x+1.0, 5.3, 1.8, 0.4, fill_color=color, radius=True)
    tag_t = add_text(s, x+1.0, 5.3, 1.8, 0.4, tag, size=11, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    animate(s, card, 'wipe_bottom', 0.5 + i*0.6, 0.6)
    animate(s, hdr, 'wipe_right', 0.8 + i*0.6, 0.3)
    animate(s, nm, 'fade', 0.9 + i*0.6, 0.4)
    animate(s, sb, 'fade', 1.0 + i*0.6, 0.3)
    animate(s, ds, 'fade', 1.1 + i*0.6, 0.5)
    animate(s, tag_bg, 'zoom', 1.4 + i*0.6, 0.3)
    animate(s, tag_t, 'fade', 1.5 + i*0.6, 0.3)

conc_bg = add_rect(s, 2, 6.3, 9.33, 0.6, fill_color=C['card'], border_color=C['accent'], border_width=1.5, radius=True)
conc = add_text(s, 2, 6.32, 9.33, 0.55, '不被任何一家闭源厂商锁死 — 竞争越激烈，你用的AI越便宜',
               size=14, bold=True, color=C['accent'], align=PP_ALIGN.CENTER)
animate(s, conc_bg, 'wipe_bottom', 3.5, 0.5)
animate(s, conc, 'fade', 3.8, 0.5)
add_footer(s, '来源：GPUStack · 知乎 · 华为昇腾公告')
flush_animations(s)

# ===== Slide 23: 章节 - 我的判断 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
add_section_page(s, '06', '我的判断', '中国团队可以在最前沿的底层机制上做原创')
flush_animations(s)

# ===== Slide 24: 结尾 =====
s = prs.slides.add_slide(blank)
set_bg(s)
add_slide_transition(s)
shape_line = add_rect(s, 4.17, 2.5, 5, 0.04, fill_color=C['accent'])
shape_main = add_text(s, 1.5, 2.8, 10.33, 1.2, '关注 AI大家学', size=40, bold=True, color=C['primary'], align=PP_ALIGN.CENTER)
shape_sub = add_text(s, 1.5, 4.0, 10.33, 0.6, 'AI 世界快人一步', size=18, color=C['subtext'], align=PP_ALIGN.CENTER)
shape_cta = add_text(s, 1.5, 5.2, 10.33, 0.5, '开源 vs 闭源，你站哪边？评论区聊聊 👇', size=13, color=C['accent'], align=PP_ALIGN.CENTER)
shape_update = add_text(s, 1.5, 5.8, 10.33, 0.4, '港股上市 · 美国调查 · 开源生态 — 三条线持续跟进', size=11, color=C['subtext'], align=PP_ALIGN.CENTER)
animate(s, shape_line, 'wipe_right', 0, 0.5)
animate(s, shape_main, 'zoom', 0.3, 0.8)
animate(s, shape_sub, 'fade', 1.0, 0.6)
animate(s, shape_cta, 'fade', 1.5, 0.5)
animate(s, shape_update, 'fade', 2.0, 0.4)
flush_animations(s)

# ========== 保存 ==========
out = os.path.join(os.path.dirname(__file__), 'KimiK3-演示文稿-v5.pptx')
prs.save(out)
print(f'PPT 已生成: {out}')
print(f'共 {len(prs.slides)} 页幻灯片')
